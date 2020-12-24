# coding=utf-8
import os
import copy
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


# 保存slide模板文本格式
def save_template_font(pres, index):
    template_text_frame = pres.slides[0].shapes[index]
    template_font = template_text_frame.text_frame.paragraphs[0].runs[0].font
    template_font_name = template_font.name
    template_font_size = int(template_font.size)
    template_font_color = template_font.color.rgb
    template_font_bold = template_font.bold
    return (template_font_name, template_font_size, template_font_color, template_font_bold)


# 修改slide文本格式
def modify_slide_font(slide, index, text_name, font_list):
    slide.shapes[index].text = text_name
    slide.shapes[index].text_frame.paragraphs[0].runs[0].font.name = font_list[0]
    slide.shapes[index].text_frame.paragraphs[0].runs[0].font.size = font_list[1]
    slide.shapes[index].text_frame.paragraphs[0].runs[0].font.color.rgb = font_list[2]
    slide.shapes[index].text_frame.paragraphs[0].runs[0].font.bold = font_list[3]
    slide.shapes[index].text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# 复制幻灯片
def duplicate_slide(pres, index):
    template = pres.slides[index]
    blank_slide_layout = pres.slide_layouts[index]
    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in dict.items(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                                    value._target,
                                                    value.rId)

    return copied_slide


# 按照slide删除幻灯片
def delete_slide_by_slide(pres, slide):
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(pres.slides._sldIdLst)}
    slide_id = slide.slide_id
    pres.part.drop_rel(id_dict[slide_id][1])
    del pres.slides._sldIdLst[id_dict[slide_id][0]]


# 按照slide索引删除幻灯片
def delete_slide_by_index(pres, index):
    rId = pres.slides._sldIdLst[index + 1].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[index + 1]


if __name__ == '__main__':
    INPUT_PATH = os.path.join(os.getcwd(), 'input')
    OUTPUT_PATH = os.path.join(os.getcwd(), 'output')
    if not os.path.exists(INPUT_PATH):
        print("存放奖项名单的以下目录不存在,请创建并放入奖项名单.\n{}".format(INPUT_PATH))
        print("奖项名单格式要求:\n"
              "1.放入{}目录下的文件夹里,比如low和high,表示使用不同的PPT模板制作;\n"
              "2.文件名称为奖项名称,文件后缀为.txt,比如 核心馒头团队.txt\n"
              "3.txt内容里填写两列,格式:获奖人员名字|获奖奖项|,比如 小汤哥|暖心后勤部长|".format(INPUT_PATH))
        exit(1)
    if not os.path.exists(OUTPUT_PATH):
        os.makedirs(OUTPUT_PATH)

    # 遍历获奖名单
    for template in os.listdir(INPUT_PATH):
        for file in os.listdir(os.path.join(INPUT_PATH, template)):
            # 区分不同的PPT模板
            ppt_template = template
            ppt = Presentation("{}.pptx".format(ppt_template))
            # 保存模板文本格式
            award_name_template_font = save_template_font(ppt, 1)
            trophy_name_template_font = save_template_font(ppt, 2)

            # 读取获奖名单
            file_name = os.path.join(INPUT_PATH, template, file)
            with open(file_name, 'r', encoding='gbk') as f:
                for line in f.readlines():
                    line_list = line.strip().split('|')
                    award_name, trophy_name = line_list[0], line_list[1]

                    # 按照获奖名单复制幻灯片
                    new_slide = duplicate_slide(ppt, 0)
                    # 修改获奖人员名字
                    modify_slide_font(new_slide, 3, award_name, award_name_template_font)
                    # 修改获奖奖项
                    modify_slide_font(new_slide, 4, trophy_name, trophy_name_template_font)

            # 删除第一张模板幻灯片
            # delete_slide_by_slide(ppt, ppt.slides[0])
            # 生成PPT保存
            ppt.save(os.path.join(OUTPUT_PATH, "{}.pptx".format(os.path.splitext(file)[0])))
            print("已生成PPT文件: {}".format(os.path.join(OUTPUT_PATH, "{}.pptx".format(os.path.splitext(file)[0]))))
